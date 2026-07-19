"""Genera la presentación editable del proyecto Mundial Trivia Challenge."""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


RAIZ = Path(__file__).resolve().parents[1]
SALIDA = RAIZ / "entregables" / "Presentacion_Trivia_Mundial.pptx"
FIGURAS = RAIZ / "docs" / "informe" / "figuras"
GRAFICO_SIMULACION = RAIZ / "resultados" / "simulacion" / "puntajes_simulados.png"

ANCHO = 13.333
ALTO = 7.5
AZUL = RGBColor(15, 37, 87)
AZUL_MEDIO = RGBColor(30, 58, 138)
AZUL_CLARO = RGBColor(59, 130, 246)
DORADO = RGBColor(251, 191, 36)
BLANCO = RGBColor(255, 255, 255)
GRIS = RGBColor(203, 213, 225)
VERDE = RGBColor(34, 197, 94)
ROJO = RGBColor(239, 68, 68)


def agregar_fondo(slide, color=AZUL):
    forma = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(ANCHO), Inches(ALTO))
    forma.fill.solid()
    forma.fill.fore_color.rgb = color
    forma.line.fill.background()
    slide.shapes._spTree.remove(forma._element)
    slide.shapes._spTree.insert(2, forma._element)


def agregar_texto(slide, texto, izquierda, arriba, ancho, alto, tamano=22,
                  color=BLANCO, negrita=False, alineacion=PP_ALIGN.LEFT,
                  fuente="Aptos"):
    caja = slide.shapes.add_textbox(
        Inches(izquierda), Inches(arriba), Inches(ancho), Inches(alto))
    marco = caja.text_frame
    marco.clear()
    marco.word_wrap = True
    marco.vertical_anchor = MSO_ANCHOR.MIDDLE
    parrafo = marco.paragraphs[0]
    parrafo.text = texto
    parrafo.alignment = alineacion
    parrafo.font.name = fuente
    parrafo.font.size = Pt(tamano)
    parrafo.font.bold = negrita
    parrafo.font.color.rgb = color
    return caja


def agregar_titulo(slide, titulo, numero):
    agregar_texto(slide, titulo, 0.65, 0.32, 11.6, 0.65, 28, DORADO, True)
    linea = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.05), Inches(12.0), Inches(0.04))
    linea.fill.solid()
    linea.fill.fore_color.rgb = AZUL_CLARO
    linea.line.fill.background()
    agregar_texto(slide, f"{numero:02d}", 12.2, 0.32, 0.45, 0.5, 14, GRIS, True,
                  PP_ALIGN.RIGHT)


def agregar_tarjeta(slide, titulo, texto, izquierda, arriba, ancho, alto,
                    acento=AZUL_CLARO):
    forma = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(izquierda), Inches(arriba), Inches(ancho), Inches(alto))
    forma.fill.solid()
    forma.fill.fore_color.rgb = RGBColor(24, 51, 108)
    forma.line.color.rgb = acento
    forma.line.width = Pt(1.5)
    agregar_texto(slide, titulo, izquierda + 0.25, arriba + 0.15,
                  ancho - 0.5, 0.45, 18, acento, True)
    agregar_texto(slide, texto, izquierda + 0.25, arriba + 0.65,
                  ancho - 0.5, alto - 0.8, 15, BLANCO)


def agregar_imagen_ajustada(slide, ruta, izquierda, arriba, ancho, alto):
    ruta = Path(ruta)
    if not ruta.exists():
        agregar_texto(slide, f"Imagen pendiente:\n{ruta.name}", izquierda, arriba,
                      ancho, alto, 16, GRIS, False, PP_ALIGN.CENTER)
        return
    with Image.open(ruta) as imagen:
        proporcion = min(ancho / imagen.width, alto / imagen.height)
        ancho_final = imagen.width * proporcion
        alto_final = imagen.height * proporcion
    x = izquierda + (ancho - ancho_final) / 2
    y = arriba + (alto - alto_final) / 2
    slide.shapes.add_picture(
        str(ruta), Inches(x), Inches(y), Inches(ancho_final), Inches(alto_final))


def agregar_notas(slide, texto):
    slide.notes_slide.notes_text_frame.text = texto


def nueva_diapositiva(presentacion, titulo, numero):
    slide = presentacion.slides.add_slide(presentacion.slide_layouts[6])
    agregar_fondo(slide)
    agregar_titulo(slide, titulo, numero)
    return slide


def generar_presentacion():
    presentacion = Presentation()
    presentacion.slide_width = Inches(ANCHO)
    presentacion.slide_height = Inches(ALTO)

    slide = presentacion.slides.add_slide(presentacion.slide_layouts[6])
    agregar_fondo(slide)
    agregar_texto(slide, "MUNDIAL", 0.75, 0.65, 5.8, 0.75, 40, DORADO, True)
    agregar_texto(slide, "TRIVIA CHALLENGE", 0.75, 1.45, 8.5, 0.9, 34, BLANCO, True)
    agregar_texto(slide, "Juego interactivo, análisis NumPy y visualización Matplotlib",
                  0.8, 2.55, 7.2, 0.9, 21, GRIS)
    agregar_texto(slide,
                  "Luis Enrique Mamani Aguilar  |  U23259985\n"
                  "Eduardo Franco Cortez Benites  |  U1421099\n"
                  "Paul Williams Gallardo Villa  |  U1614474",
                  0.8, 4.25, 6.4, 1.25, 16, BLANCO)
    agregar_imagen_ajustada(slide, FIGURAS / "03_pregunta.png", 8.2, 1.25, 4.5, 4.9)
    agregar_texto(slide,
                  "Universidad Tecnológica del Perú\nProgramación de Computadoras",
                  0.8, 6.25, 7, 0.7, 14, DORADO, True)
    agregar_notas(slide, "Presentar al equipo y explicar en una frase el propósito del proyecto. Tiempo sugerido: 35 segundos.")

    slide = nueva_diapositiva(presentacion, "Problema y solución", 2)
    agregar_tarjeta(slide, "PROBLEMA",
                    "Un puntaje final aislado no explica dónde falla el jugador, cómo evoluciona ni cuánto tarda en responder.",
                    0.75, 1.45, 5.75, 4.6, ROJO)
    agregar_tarjeta(slide, "SOLUCIÓN",
                    "Registrar cada interacción y convertirla en matriz, indicadores, clasificación, gráficos y reportes reproducibles.",
                    6.85, 1.45, 5.75, 4.6, VERDE)
    agregar_notas(slide, "Luis: explicar la situación problemática y conectar el juego con el análisis de datos. Tiempo: 45 segundos.")

    slide = nueva_diapositiva(presentacion, "Flujo de la partida", 3)
    fases = (
        ("GRUPOS", "+10", 0.7), ("OCTAVOS", "+15", 3.15),
        ("CUARTOS", "+20", 5.6), ("SEMIFINAL", "+30", 8.05),
        ("FINAL", "+50", 10.5),
    )
    for nombre, puntos, x in fases:
        agregar_tarjeta(slide, nombre, f"5 preguntas\n{puntos} puntos por acierto",
                        x, 2.0, 2.15, 2.35, DORADO)
    agregar_texto(slide, "3 vidas · 25 preguntas máximas · sin repeticiones",
                  2.0, 5.2, 9.3, 0.7, 24, BLANCO, True, PP_ALIGN.CENTER)
    agregar_notas(slide, "Eduardo: describir fases, puntajes, vidas y condición de finalización. Tiempo: 45 segundos.")

    slide = nueva_diapositiva(presentacion, "Arquitectura modular", 4)
    modulos = (
        ("interfaz.py", "Pantallas y eventos", 0.8, 1.55),
        ("juego.py", "Reglas y validaciones", 4.65, 1.55),
        ("preguntas.py", "50 preguntas", 8.5, 1.55),
        ("analitica.py", "NumPy y Matplotlib", 2.7, 4.2),
        ("simulador.py", "Partidas reproducibles", 6.55, 4.2),
    )
    for nombre, descripcion, x, y in modulos:
        agregar_tarjeta(slide, nombre, descripcion, x, y, 3.1, 1.45,
                        DORADO if "analitica" in nombre else AZUL_CLARO)
    agregar_notas(slide, "Paul: explicar que la interfaz no contiene las reglas y que la misma lógica se reutiliza en pruebas y simulación. Tiempo: 50 segundos.")

    slide = nueva_diapositiva(presentacion, "Estructuras de datos con propósito", 5)
    agregar_tarjeta(slide, "LISTAS", "Banco de preguntas\nHistorial\nIDs respondidos\nRanking",
                    0.75, 1.45, 3.75, 4.9, AZUL_CLARO)
    agregar_tarjeta(slide, "DICCIONARIOS", "Pregunta\nEstado del juego\nRegistro por respuesta\nCódigos de categorías",
                    4.8, 1.45, 3.75, 4.9, DORADO)
    agregar_tarjeta(slide, "TUPLAS", "Opciones inmutables\nFases del torneo\nColumnas de matrices\nCampos CSV",
                    8.85, 1.45, 3.75, 4.9, VERDE)
    agregar_notas(slide, "Luis: justificar por qué se eligió cada estructura, no limitarse a nombrarlas. Tiempo: 55 segundos.")

    slide = nueva_diapositiva(presentacion, "Captura y matriz NumPy", 6)
    columnas = (
        "Número", "Pregunta", "Fase", "Categoría", "Dificultad",
        "Acierto", "Puntos", "Puntaje", "Vidas", "Tiempo",
    )
    for indice, columna in enumerate(columnas):
        fila = indice // 5
        columna_pos = indice % 5
        x = 0.75 + columna_pos * 2.48
        y = 1.55 + fila * 1.35
        agregar_tarjeta(slide, str(indice), columna, x, y, 2.15, 1.0,
                        DORADO if indice in (5, 7, 9) else AZUL_CLARO)
    agregar_texto(slide,
                  "Matriz n × 10  →  sumas, medias, desviación estándar y máscaras booleanas",
                  1.1, 4.75, 11.1, 0.85, 22, BLANCO, True, PP_ALIGN.CENTER)
    agregar_notas(slide, "Eduardo: mostrar cómo un diccionario se transforma en fila numérica y nombrar las operaciones NumPy. Tiempo: 60 segundos.")

    slide = nueva_diapositiva(presentacion, "Clasificación y visualización", 7)
    clasificaciones = (
        ("Excelente", "80–100%", VERDE), ("Bueno", "60–79.99%", AZUL_CLARO),
        ("En proceso", "40–59.99%", DORADO), ("Necesita refuerzo", "0–39.99%", ROJO),
    )
    for indice, (nombre, rango, color) in enumerate(clasificaciones):
        agregar_tarjeta(slide, nombre, rango, 0.8, 1.4 + indice * 1.25, 4.2, 0.95, color)
    agregar_imagen_ajustada(slide, GRAFICO_SIMULACION, 5.35, 1.35, 7.3, 5.3)
    agregar_notas(slide, "Paul: explicar los límites de clasificación y presentar el histograma. Tiempo: 55 segundos.")

    slide = nueva_diapositiva(presentacion, "Simulación reproducible", 8)
    metricas = (
        ("10", "partidas"), ("104.50", "puntaje promedio"),
        ("175", "puntaje máximo"), ("43.27", "desviación estándar"),
    )
    for indice, (valor, etiqueta) in enumerate(metricas):
        x = 0.8 + (indice % 2) * 3.25
        y = 1.5 + (indice // 2) * 2.0
        agregar_tarjeta(slide, valor, etiqueta, x, y, 2.85, 1.55,
                        DORADO if indice == 1 else AZUL_CLARO)
    agregar_texto(slide, "Semilla 42", 8.1, 1.75, 3.8, 0.7, 32, DORADO, True,
                  PP_ALIGN.CENTER)
    agregar_texto(slide,
                  "Misma semilla\n=\nmismas preguntas y resultados",
                  8.1, 2.75, 3.8, 2.0, 22, BLANCO, True, PP_ALIGN.CENTER)
    agregar_notas(slide, "Eduardo: ejecutar el simulador o mostrar sus archivos; aclarar que la semilla vuelve repetible el escenario. Tiempo: 55 segundos.")

    slide = nueva_diapositiva(presentacion, "Validaciones y pruebas", 9)
    agregar_tarjeta(slide, "VALIDACIONES",
                    "Nombre y longitud\nBanco consistente\nIDs únicos\nCuatro opciones\nPorcentajes válidos\nSin respuestas duplicadas",
                    0.8, 1.45, 5.7, 4.9, DORADO)
    agregar_tarjeta(slide, "11 PRUEBAS SUPERADAS",
                    "Lógica del juego\nMatriz e indicadores\nLímites de clasificación\nExportación de gráficos\nRepetibilidad\nArchivos de simulación",
                    6.85, 1.45, 5.7, 4.9, VERDE)
    agregar_notas(slide, "Luis: mostrar el comando unittest y explicar dos casos de frontera. Tiempo: 50 segundos.")

    slide = nueva_diapositiva(presentacion, "Evidencia de funcionamiento", 10)
    agregar_imagen_ajustada(slide, FIGURAS / "01_bienvenida.png", 0.65, 1.35, 6.0, 4.8)
    agregar_imagen_ajustada(slide, FIGURAS / "05_resultado.png", 6.75, 1.35, 6.0, 4.8)
    agregar_texto(slide, "Entrada validada", 1.6, 6.2, 4.0, 0.4, 16, DORADO, True,
                  PP_ALIGN.CENTER)
    agregar_texto(slide, "Reporte y clasificación", 7.75, 6.2, 4.0, 0.4, 16, DORADO, True,
                  PP_ALIGN.CENTER)
    agregar_notas(slide, "Paul: realizar la demostración en vivo desde el ingreso del nombre hasta una respuesta y mostrar la carpeta resultados. Tiempo: 90 segundos.")

    slide = nueva_diapositiva(presentacion, "Entregables reproducibles", 11)
    agregar_tarjeta(slide, "PROGRAMA", "Python\nSimulador\nPruebas\nImágenes",
                    0.75, 1.5, 3.75, 3.8, AZUL_CLARO)
    agregar_tarjeta(slide, "DOCUMENTACIÓN", "LaTeX APA 7\nPDF\nWord\nReferencias",
                    4.8, 1.5, 3.75, 3.8, DORADO)
    agregar_tarjeta(slide, "SUSTENTACIÓN", "PPT editable\nCapturas\nGráficos\nGuion de video",
                    8.85, 1.5, 3.75, 3.8, VERDE)
    agregar_texto(slide, "docker compose run --rm latex  |  docker compose run --rm word",
                  1.0, 5.75, 11.3, 0.75, 19, BLANCO, True, PP_ALIGN.CENTER)
    agregar_notas(slide, "Explicar que PDF y Word se regeneran con Docker sin instalar LaTeX localmente. Tiempo: 45 segundos.")

    slide = nueva_diapositiva(presentacion, "Conclusiones", 12)
    agregar_texto(slide,
                  "DATOS REALES\nCada respuesta queda registrada",
                  0.8, 1.6, 3.7, 1.4, 23, BLANCO, True, PP_ALIGN.CENTER)
    agregar_texto(slide,
                  "ANÁLISIS ÚTIL\nNumPy y gráficos interpretables",
                  4.8, 1.6, 3.7, 1.4, 23, BLANCO, True, PP_ALIGN.CENTER)
    agregar_texto(slide,
                  "CALIDAD\nValidaciones, pruebas y Docker",
                  8.8, 1.6, 3.7, 1.4, 23, BLANCO, True, PP_ALIGN.CENTER)
    agregar_texto(slide, "¿Preguntas?", 2.4, 4.4, 8.5, 1.1, 42, DORADO, True,
                  PP_ALIGN.CENTER)
    agregar_notas(slide, "Cerrar relacionando los resultados con la rúbrica y abrir el espacio de preguntas. Tiempo: 30 segundos.")

    SALIDA.parent.mkdir(parents=True, exist_ok=True)
    presentacion.save(SALIDA)
    print(f"Presentación creada: {SALIDA}")


if __name__ == "__main__":
    generar_presentacion()
